import numpy as np
import pandas as pd
from math import *
import six

from PandasToPowerpoint import df_to_table
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE

pd.options.display.max_columns = 999

round_to_n = lambda x, n: round(x, -int(floor(log10(abs(x)))) + (n - 1))

def _do_formatting(value, format_str):
    """Format value according to format_str, and deal
    sensibly with format_str if it is missing or invalid."""
    if format_str == '':
        if type(value) in six.integer_types:
            format_str = ','
        elif type(value) is float:
            format_str = 'f'
        elif type(value) is str:
            format_str = 's'
    elif format_str[0] == '.':
        if format_str.endswith('R'):
            if type(value) in six.integer_types:
                value = round_to_n(value, int(format_str[1]))
                format_str = ','
        if not format_str.endswith('G'):
            format_str = format_str + "G"
    try:
        value = format(value, format_str)
    except:
        value = format(value, '')

    return value

def process_position_parameter(param):
    """Process positioning parameters (left, top, width, height) given to
    df_to_table.
    If an integer, returns the right instance of the Cm class to allow it to be treated
    as cm. If missing, then default to 4cm. Otherwise, pass through whatever it gets.
    """
    if param is None:
        return Cm(4)
    elif type(param) is int:
        return Cm(param)
    else:
        return param


def df_to_table(slide, df, left=None, top=None, width=None, height=None,
                colnames=None, rownames=None, col_formatters=None, rounding=None,
                name=None, white_backgr=False, transposed=False, font_size=14, set_col_width=False,
                col_width=[], set_col_height=False, col_height=[]):

    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.

    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.

    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data

    Optional arguments:
     - left: Position of the left-side of the table, either as an integer in cm, or as an instance of a
     pptx.util Length class (pptx.util.Inches for example). Defaults to 4cm.
     - top: Position of the top of the table, takes parameters as above.
     - width: Width of the table, takes parameters as above.
     - height: Height of the table, takes parameters as above.
     - colnames: Column names
     - rownames: Row names
     - col_formatters: A n_columns element long list containing format specifications for each column.
     For example ['', ',', '.2'] does no special formatting for the first column, uses commas as thousands separators
     in the second column, and formats the third column as a float with 2 decimal places.
     - rounding: A n_columns element long list containing a number for each integer column that requires rounding
     that is then multiplied by -1 and passed to round().
     - name: A name to be given to the table in the Powerpoint file. This is not displayed, but can help
     extract the table later to make further changes.
     - white_backgr: Set False for customizing cell background color, font etc
     - transposed: Set True if the table has horizontal orientation.
     - font_size: font size of table labels
     - set_col_width: Set True for manually setting column width
     - col_width: List of column widths. Needs to have as many elements as the columns of the table. For example, for a table
      with two columns it can be: col_width=[Inches(1.3),Inches(1.2)]
     - set_col_height: Assign True for manually setting column width
     - col_height: List of column heights. Needs to have as many elements as the columns of the table. For example, for a table
      with two columns it can be: col_height=[Inches(1.3),Inches(1.2)]
     """

    left = process_position_parameter(left)
    top = process_position_parameter(top)
    width = process_position_parameter(width)
    height = process_position_parameter(height)

    rows, cols = df.shape

    # if table has horizontal orientation
    if transposed:
        shp = slide.shapes.add_table(rows, cols+1, left, top, width, height)

        if rownames is None:
            rownames = list(df.index)

        # Insert the row names
        for row_index, row_name in enumerate(rownames):
            shp.table.cell(row_index,0).text = row_name

        m = df.as_matrix()

        for row in range(rows):
            for col in range(cols):
                val = m[row, col]
                if col_formatters is None:
                    text = str(val)
                else:
                    text = _do_formatting(val, col_formatters[col])
                shp.table.cell(row, col+1).text = text

        if name is not None:
            shp.name = name

        # format text within cells
        if not white_backgr:
            for i in range(rows):
                for j in range(cols+1):
                    shp.table.cell(i, j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    para = shp.table.cell(i, j).text_frame.paragraphs[0]
                    para.font.bold = False
                    para.font.size = Pt(font_size)
                    para.font.name = 'Calibri'
                    para.font.color.rgb = RGBColor(0, 0, 0)
                    para.alignment = PP_ALIGN.CENTER

            # format labels
            for i in range(rows):
                    shp.table.cell(i, 0).fill.solid()
                    shp.table.cell(i, 0).fill.fore_color.rgb = RGBColor(79, 129, 189)
                    para = shp.table.cell(i, 0).text_frame.paragraphs[0]
                    para.font.bold = True
                    para.font.size = Pt(font_size)
                    para.font.name = 'Calibri'
                    para.font.color.rgb = RGBColor(0, 0, 0)
                    para.alignment = PP_ALIGN.LEFT

            for j in range(cols):
                    shp.table.cell(0, j+1).fill.solid()
                    shp.table.cell(0, j+1).fill.fore_color.rgb = RGBColor(233, 237, 244)

    # if table has vertical orientation
    else:
        shp = slide.shapes.add_table(rows+1, cols, left, top, width, height)

        if colnames is None:
            colnames = list(df.columns)

        # Insert the column names
        for col_index, col_name in enumerate(colnames):
            shp.table.cell(0,col_index).text = col_name

        m = df.as_matrix()

        for row in range(rows):
            for col in range(cols):
                val = m[row, col]

                if col_formatters is None:
                    text = str(val)
                else:
                    text = _do_formatting(val, col_formatters[col])

                shp.table.cell(row+1, col).text = text

        if name is not None:
            shp.name = name

        if not white_backgr:
            for i in range(rows+1):
                for j in range(cols):
                    shp.table.cell(i, j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    para = shp.table.cell(i, j).text_frame.paragraphs[0]
                    para.font.bold = False
                    para.font.size = Pt(font_size)
                    para.font.name = 'Calibri'
                    para.font.color.rgb = RGBColor(0, 0, 0)
                    para.alignment = PP_ALIGN.CENTER

            # format labels
            for j in range(cols):
                    para = shp.table.cell(0, j).text_frame.paragraphs[0]
                    para.font.bold = True
                    para.font.size = Pt(font_size)
                    para.font.name = 'Calibri'
                    para.font.color.rgb = RGBColor(0, 0, 0)
                    para.alignment = PP_ALIGN.CENTER

    if set_col_width:
        for i in range(len(col_width)):
            shp.table.columns[i].width = col_width[i]

    if set_col_height:
        for i in range(len(col_height)):
            shp.table.columns[i].height = col_height[i]

    return shp


def add_slide():
    title_only_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    return slide


def format_table(table, rows, columns):
    for i in range(rows):
        for j in range(columns):
            _set_cell_border(table.cell(i, j))
            para = table.cell(i, j).text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(0, 0, 0)
            para.alignment = PP_ALIGN.CENTER
        for i in range(rows):
            for j in range(columns):
                table.cell(i, j).fill.background()


def add_title(slide, left, top, width, height, title, bold=True, font_size=16, underline=True, italics=False):
    txBox = slide.shapes.add_textbox(left=left, top=top, width=width, height=height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = bold
    p.font.underline = underline
    p.font.italics = italics
    p.font.size = Pt(font_size)
    p.alignment = PP_ALIGN.CENTER


# Example run:

# prs = Presentation()
# slide = add_slide()
# add_title(slide,left=Inches(0.4), top=Inches(1.4), width=Inches(2.2), height=Inches(0.2), title='2017 Net Sales',
#                  font_size=11)

# df_to_table(slide, salestable, left=Inches(0.1), top=Inches(2.0), width=Inches(3.0), height=Inches(0.4),colnames=None,
#                    col_formatters=None, rounding=None, name=None, transposed=True, font_size=11,
#                    set_col_width=True, col_width=[Inches(1.3),Inches(1.2)])
